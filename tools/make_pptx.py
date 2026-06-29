"""生成 MAOE 路演 .pptx —— 16:9 高级科技感深蓝主题，8 页。

使用：
    .venv/bin/python tools/make_pptx.py
输出：
    docs/slides/MAOE-roadshow.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ============== 配色 ==============
BG = RGBColor(0x08, 0x0C, 0x1C)
BG2 = RGBColor(0x12, 0x1C, 0x38)
ACCENT = RGBColor(0x58, 0xC6, 0xFF)        # 青
ACCENT_2 = RGBColor(0xFF, 0xB8, 0x4D)      # 金
TEXT = RGBColor(0xEB, 0xF0, 0xFA)
TEXT_DIM = RGBColor(0xA0, 0xB4, 0xD2)
NODE_LINE = RGBColor(0x3C, 0x5E, 0xA0)

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "MAOE-roadshow.pptx"


def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # 把背景放最底
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _accent_bar(slide, x, y, w=Inches(2.4), h=Inches(0.18)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT


def _text(slide, x, y, w, h, text: str, size: int, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = "PingFang SC"
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
    return tb


def _node(slide, x, y, w, h, text: str, accent=ACCENT, font_size=18):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.line.color.rgb = accent
    box.line.width = Pt(2)
    box.fill.solid()
    box.fill.fore_color.rgb = BG2
    # 左色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), h)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    # 文字
    tb = box.text_frame
    tb.margin_left = Inches(0.2)
    tb.margin_top = Inches(0.05)
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "PingFang SC"
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = TEXT
    return box


def _arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(2, x1, y1, x2, y2)
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2.5)


def _header(slide, page_num: int, total: int, subtitle: str):
    # 页码
    _text(slide, Inches(12.0), Inches(0.25), Inches(1.2), Inches(0.4),
          f"{page_num:02d} / {total:02d}", 12, color=TEXT_DIM, align=PP_ALIGN.RIGHT)
    # 顶部条
    _accent_bar(slide, Inches(0.6), Inches(0.55))
    _text(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.5),
          subtitle, 16, color=ACCENT_2)


# ============== 各页 ==============
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, BG)
    _text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(2),
          "MAOE", 130, color=TEXT, bold=True)
    # 右侧色条
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.6), Inches(2.4),
                              Inches(0.16), Inches(1.8))
    bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    _text(s, Inches(0.8), Inches(4.0), Inches(11), Inches(1),
          "Multi-Agent Dynamic Orchestration Engine", 32, color=ACCENT, bold=True)
    _text(s, Inches(0.8), Inches(4.8), Inches(12), Inches(1),
          "Skill 原生工作流编译 · 三档候选 Pareto 调度 · 决策可追溯", 20,
          color=TEXT_DIM)
    _text(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
          "海聚英才 OPC 专项赛 · 极客组参赛项目", 18, color=ACCENT_2, bold=True)
    _text(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
          "github.com/zhongxiaomi06-sudo/maoe", 14, color=TEXT_DIM)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, BG)
    _header(s, 2, 8, "01 我们要解决的问题")
    _text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.8),
          "Agent 工作流的三个痛点", 36, color=TEXT, bold=True)

    items = [
        ("拍脑袋编排", "工作流靠人工试错，谁也说不清为什么选这个 skill。"),
        ("成本不可预算", "跑起来才知道烧多少 token，预算控制全靠 hardcode 限流。"),
        ("决策不可追溯", "为什么这次选了 gpt-4o-mini？没有日志可查、没法回放。"),
    ]
    for i, (t, b) in enumerate(items):
        y = Inches(2.7 + i * 1.3)
        # 序号方块
        num = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y,
                                  Inches(0.9), Inches(0.9))
        num.line.fill.background(); num.fill.solid(); num.fill.fore_color.rgb = ACCENT
        tf = num.text_frame
        tf.margin_left = tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"0{i+1}"
        r.font.name = "PingFang SC"; r.font.size = Pt(28)
        r.font.bold = True; r.font.color.rgb = BG
        _text(s, Inches(1.7), y, Inches(11), Inches(0.5),
              t, 24, color=ACCENT_2, bold=True)
        _text(s, Inches(1.7), y + Inches(0.5), Inches(11), Inches(0.5),
              b, 16, color=TEXT_DIM)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, BG)
    _header(s, 3, 8, "02 MAOE 的回答")
    _text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.8),
          "Skill 即类型契约 + 编译期决策", 36, color=TEXT, bold=True)

    # 流程图：goal → IR → graph → 三档候选 → pareto → run
    nodes = [
        ("自然语言 goal", Inches(0.4), 1.6),
        ("GoalIR", Inches(2.4), 1.6),
        ("CapabilityGraph", Inches(2.4), 1.6),
        ("3 档 WorkflowDAG", Inches(2.4), 1.6),
        ("Pareto 选档", Inches(2.0), 1.6),
        ("Runtime 落盘", Inches(2.0), 1.6),
    ]
    y = Inches(3.2)
    x = Inches(0.4)
    boxes = []
    for label, w, h in nodes:
        _node(s, x, y, w, Inches(h), label, accent=ACCENT, font_size=14)
        boxes.append((x, y, w))
        x += w + Inches(0.15)
    # 箭头
    for i in range(len(boxes) - 1):
        x0, y0, w0 = boxes[i]
        x1, y1, _ = boxes[i + 1]
        _arrow(s, x0 + w0, y0 + Inches(0.8), x1, y1 + Inches(0.8))

    _text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5),
          "  · 编译器先做 5 类静态校验：unknown_dep / cycle / capability_gap / schema_mismatch / missing_deliverable\n"
          "  · 同一 goal 输出 economy / balanced / quality 三档候选，五维 Pareto 前沿选推荐档\n"
          "  · 编译产出 sha256 lock_digest，决策 100% 可复现",
          16, color=TEXT_DIM)


def slide_arch(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, BG)
    _header(s, 4, 8, "03 系统架构")
    _text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.8),
          "两条平行管线：运行 + 编译", 36, color=TEXT, bold=True)

    # 主轴
    cx = SLIDE_W / 2
    _node(s, cx - Inches(2.4), Inches(2.5), Inches(4.8), Inches(0.65),
          "CLI: maoe run / benchmark", ACCENT_2, 16)
    _node(s, cx - Inches(2.4), Inches(3.4), Inches(4.8), Inches(0.65),
          "MAOEEngine（运行管线总指挥）", ACCENT, 16)
    _arrow(s, cx, Inches(3.15), cx, Inches(3.4))

    titles = ["Parser", "Evaluator", "Router", "Economist", "Executor", "Quality"]
    bw, bh, gap = Inches(1.45), Inches(0.55), Inches(0.1)
    for i, t in enumerate(titles[:3]):
        x = cx - Inches(2.4) + i * (bw + gap)
        _node(s, x, Inches(4.4), bw, bh, t, ACCENT, 12)
    for i, t in enumerate(titles[3:]):
        x = cx - Inches(2.4) + i * (bw + gap)
        _node(s, x, Inches(5.1), bw, bh, t, ACCENT, 12)
    _arrow(s, cx, Inches(4.05), cx, Inches(4.4))

    _node(s, cx - Inches(2.4), Inches(6.0), Inches(4.8), Inches(0.55),
          "LLM Client → 代理 LLM API", ACCENT_2, 14)
    _arrow(s, cx, Inches(5.65), cx, Inches(6.0))

    # 旁路
    side_lefts = [
        ("Compiler", Inches(2.5)),
        ("Experiments", Inches(3.3)),
        ("Runtime", Inches(4.1)),
    ]
    for lab, y in side_lefts:
        _node(s, Inches(0.4), y, Inches(2.4), Inches(0.5), lab, NODE_LINE, 12)
    side_rights = [
        ("Registry", Inches(2.5)),
        ("Bootstrap", Inches(3.3)),
        ("Models", Inches(4.1)),
    ]
    for lab, y in side_rights:
        _node(s, SLIDE_W - Inches(2.8), y, Inches(2.4), Inches(0.5), lab, NODE_LINE, 12)


def slide_innovation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, BG)
    _header(s, 5, 8, "04 关键创新")
    _text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.8),
          "四张差异化卡片", 36, color=TEXT, bold=True)

    cards = [
        ("Skill 即类型契约",
         "YAML 胶囊声明 provides / requires / schemas / risk / rollback；\n"
         "编译期就能发现 capability gap，避免烧钱跑错。"),
        ("三档候选 Pareto",
         "对同一 goal 编译 economy/balanced/quality 三档；\n"
         "(quality↑ reliability↑ cost↓ latency↓ risk↓) 五维选最优。"),
        ("决策可复现 lockfile",
         "编译产出 sha256 lock_digest 写 workflow.lock.json；\n"
         "Runtime append-only 事件/决策流，全程可回放。"),
        ("内置 redaction 安全",
         "落盘前自动脱敏 sk- / Bearer / api_key / authorization；\n"
         "5 类 OWASP 风格静态检查兜底。"),
    ]
    for i, (t, b) in enumerate(cards):
        row, col = divmod(i, 2)
        x = Inches(0.5 + col * 6.3)
        y = Inches(2.4 + row * 2.4)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(6.0), Inches(2.1))
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.5)
        box.fill.solid(); box.fill.fore_color.rgb = BG2
        _text(s, x + Inches(0.3), y + Inches(0.2), Inches(5.5), Inches(0.5),
              f"0{i+1}  {t}", 22, color=ACCENT_2, bold=True)
        _text(s, x + Inches(0.3), y + Inches(0.95), Inches(5.5), Inches(1),
              b, 14, color=TEXT_DIM)


def slide_metrics(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, BG)
    _header(s, 6, 8, "05 工程现状")
    _text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.8),
          "可量化的工程产出", 36, color=TEXT, bold=True)

    metrics = [
        ("48", "测试用例"),
        ("100%", "ruff 干净"),
        ("10", "Skill 胶囊"),
        ("8", "Agent 角色"),
        ("3", "候选档位"),
        ("5", "静态校验"),
    ]
    cw, ch = Inches(4.0), Inches(2.2)
    gap_x, gap_y = Inches(0.2), Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(2.5)
    for i, (n, lab) in enumerate(metrics):
        row, col = divmod(i, 3)
        x = start_x + col * (cw + gap_x)
        y = start_y + row * (ch + gap_y)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, ch)
        box.line.color.rgb = ACCENT; box.line.width = Pt(2)
        box.fill.solid(); box.fill.fore_color.rgb = BG2
        _text(s, x + Inches(0.3), y + Inches(0.2), cw, Inches(1.4),
              n, 72, color=ACCENT_2, bold=True)
        _text(s, x + Inches(0.3), y + Inches(1.5), cw, Inches(0.5),
              lab, 16, color=TEXT_DIM)


def slide_deploy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, BG)
    _header(s, 7, 8, "06 演示与部署")
    _text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.8),
          "三条命令即可重现", 36, color=TEXT, bold=True)

    # 左：命令
    cmd_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5),
                                   Inches(7.2), Inches(4.5))
    cmd_box.line.color.rgb = ACCENT; cmd_box.line.width = Pt(1.5)
    cmd_box.fill.solid(); cmd_box.fill.fore_color.rgb = BG2
    _text(s, Inches(0.7), Inches(2.6), Inches(7), Inches(0.5),
          "$ ", 18, color=ACCENT, bold=True)
    cmds = [
        ("git clone github.com/zhongxiaomi06-sudo/maoe", "克隆仓库"),
        ("uv sync --extra dev", "同步依赖 (27 包)"),
        ("export MAOE_API_KEY=<your-key>", "配置鉴权"),
        (".venv/bin/maoe run \"写一个反转字符串的函数\"", "端到端跑一次"),
        (".venv/bin/pytest -q", "48/48 测试"),
        (".venv/bin/maoe benchmark -p default", "跑 benchmark"),
    ]
    y = Inches(2.95)
    for cmd, desc in cmds:
        _text(s, Inches(0.7), y, Inches(7), Inches(0.4), cmd, 14,
              color=TEXT, bold=True)
        _text(s, Inches(0.7), y + Inches(0.32), Inches(7), Inches(0.3),
              "  → " + desc, 11, color=TEXT_DIM)
        y += Inches(0.65)

    # 右：链接 + 文档
    _text(s, Inches(8.2), Inches(2.6), Inches(5), Inches(0.5),
          "开源链接", 20, color=ACCENT_2, bold=True)
    _text(s, Inches(8.2), Inches(3.1), Inches(5), Inches(0.5),
          "github.com/zhongxiaomi06-sudo/maoe", 12, color=TEXT)
    _text(s, Inches(8.2), Inches(4.0), Inches(5), Inches(0.5),
          "学习文档", 20, color=ACCENT_2, bold=True)
    _text(s, Inches(8.2), Inches(4.5), Inches(5), Inches(0.5),
          "docs/MAOE-源码逐文件讲解.md", 12, color=TEXT)
    _text(s, Inches(8.2), Inches(4.85), Inches(5), Inches(0.5),
          "  · 650+ 行逐文件讲解，含架构图与时序图", 11, color=TEXT_DIM)
    _text(s, Inches(8.2), Inches(5.7), Inches(5), Inches(0.5),
          "海报", 20, color=ACCENT_2, bold=True)
    _text(s, Inches(8.2), Inches(6.2), Inches(5), Inches(0.5),
          "docs/poster/maoe-poster-65x110.jpg", 12, color=TEXT)


def slide_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, BG)
    # 大色块
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.6),
                              Inches(0.2), Inches(4))
    bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    _text(s, Inches(1.2), Inches(1.5), Inches(11), Inches(2),
          "MAOE", 110, color=TEXT, bold=True)
    _text(s, Inches(1.2), Inches(3.5), Inches(11), Inches(1.5),
          "让 Agent 工程像编译器一样长大", 36, color=ACCENT, bold=True)
    _text(s, Inches(1.2), Inches(4.7), Inches(11), Inches(0.8),
          "Skill 即契约 · 决策即代码 · 运行即审计", 22, color=TEXT_DIM)

    _text(s, Inches(1.2), Inches(6.4), Inches(11), Inches(0.6),
          "海聚英才 OPC 极客组  ·  github.com/zhongxiaomi06-sudo/maoe", 14,
          color=TEXT_DIM)


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_arch(prs)
    slide_innovation(prs)
    slide_metrics(prs)
    slide_deploy(prs)
    slide_close(prs)

    prs.save(OUT_PATH)
    print(f"pptx saved: {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
